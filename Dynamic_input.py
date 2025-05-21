import os
from typing import List, Literal
from airflow import DAG
from airflow.decorators import task
from datetime import datetime
from pydantic import BaseModel, Field
from office365.sharepoint.client_context import ClientContext
from office365.runtime.auth.client_credential import ClientCredential
from office365.runtime.client_request_exception import ClientRequestException
from dotenv import load_dotenv
import logging
import io
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
import eml_parser
from airflow_ai_sdk.tools import tool

# Configure logging
logging.basicConfig(filename="document_router.log", level=logging.INFO)

# Load .env file and set OpenAI API key
load_dotenv()
api_key = os.getenv("OPENAI_API_KEY")
if api_key:
    os.environ["OPENAI_API_KEY"] = api_key
else:
    logging.error("OPENAI_API_KEY not found")
    raise ValueError("OPENAI_API_KEY is required")

# Pydantic models
class DocumentMetadata(BaseModel):
    path: str = Field(description="Path or URL of the document")
    content_type: str = Field(description="File extension")
    has_images: bool = Field(description="Whether the document contains images")
    is_sharepoint: bool = Field(description="Whether the document is from SharePoint")

class AgentAssignment(BaseModel):
    document: DocumentMetadata
    agent_type: Literal["file_parsing", "email", "transcript", "image_processing"]
    priority: int = Field(description="Processing priority (1-5)", ge=1, le=5)

# Image detection functions
def has_images_in_pdf(file_path: str) -> bool:
    try:
        images = convert_from_path(file_path, first_page=1, last_page=1)
        return len(images) > 0
    except Exception as e:
        logging.warning(f"Error checking images in PDF {file_path}: {str(e)}")
        return False

def has_images_in_docx(file_path: str) -> bool:
    try:
        doc = Document(file_path)
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                return True
        return False
    except Exception as e:
        logging.warning(f"Error checking images in DOCX {file_path}: {str(e)}")
        return False

def has_images_in_pptx(file_path: str) -> bool:
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    return True
        return False
    except Exception as e:
        logging.warning(f"Error checking images in PPTX {file_path}: {str(e)}")
        return False

def has_images_in_eml(file_path: str) -> bool:
    try:
        with open(file_path, "r") as f:
            eml = eml_parser.EmlParser().decode_email(f.read())
            for attachment in eml.get("attachment", []):
                if attachment["content_type"].startswith("image/"):
                    return True
            return False
    except Exception as e:
        logging.warning(f"Error checking images in EML {file_path}: {str(e)}")
        return False

@task
def validate_and_extract_metadata(source: str, is_sharepoint: bool = False) -> DocumentMetadata:
    client_id = os.getenv("SHAREPOINT_CLIENT_ID")
    client_secret = os.getenv("SHAREPOINT_CLIENT_SECRET")
    tenant_id = os.getenv("SHAREPOINT_TENANT_ID")
    if not all([client_id, client_secret, tenant_id]) and is_sharepoint:
        logging.error("Missing SharePoint credentials")
        raise ValueError("Missing SharePoint credentials")

    site_url = "[invalid url, do not cite]

    if is_sharepoint:
        try:
            ctx = ClientContext(site_url).with_credentials(ClientCredential(client_id, client_secret))
            if "AllItems.aspx" in source:
                folder_path = "/teams/tensaiGPT-PROD-HR-Docs/Shared Documents"
                folder = ctx.web.get_folder_by_server_relative_url(folder_path)
                files = folder.files
                ctx.load(files)
                ctx.execute_query()
                # Process first file for simplicity
                if not files:
                    raise ValueError(f"No files found in folder {folder_path}")
                file = files[0]
                content = io.BytesIO()
                file.download(content).execute_query()
                content.seek(0)
                temp_file = f"/tmp/{file.properties['Name']}"
                with open(temp_file, "wb") as f:
                    f.write(content.read())
                content_type = file.properties["Name"].split(".")[-1].lower()
                has_images = False
                if content_type == "pdf":
                    has_images = has_images_in_pdf(temp_file)
                elif content_type == "docx":
                    has_images = has_images_in_docx(temp_file)
                elif content_type in ["ppt", "pptx"]:
                    has_images = has_images_in_pptx(temp_file)
                elif content_type == "eml":
                    has_images = has_images_in_eml(temp_file)
                if os.path.exists(temp_file):
                    os.remove(temp_file)
                return DocumentMetadata(
                    path=source,
                    content_type=content_type,
                    has_images=has_images,
                    is_sharepoint=True
                )
            else:
                relative_url = source.split("GetFileByServerRelativeUrl")[1][1:-1]
                file = ctx.web.get_file_by_server_relative_url(relative_url).execute_query()
                content = io.BytesIO()
                file.download(content).execute_query()
                content.seek(0)
                temp_file = f"/tmp/{file.properties['Name']}"
                with open(temp_file, "wb") as f:
                    f.write(content.read())
                content_type = file.properties["Name"].split(".")[-1].lower()
                has_images = False
                if content_type == "pdf":
                    has_images = has_images_in_pdf(temp_file)
                elif content_type == "docx":
                    has_images = has_images_in_docx(temp_file)
                elif content_type in ["ppt", "pptx"]:
                    has_images = has_images_in_pptx(temp_file)
                elif content_type == "eml":
                    has_images = has_images_in_eml(temp_file)
                if os.path.exists(temp_file):
                    os.remove(temp_file)
                return DocumentMetadata(
                    path=source,
                    content_type=content_type,
                    has_images=has_images,
                    is_sharepoint=True
                )
        except Exception as e:
            logging.error(f"SharePoint validation failed for {source}: {str(e)}")
            raise
    else:
        if not os.path.exists(source):
            logging.error(f"Local file not found: {source}")
            raise ValueError(f"Local file not found: {source}")
        content_type = source.split(".")[-1].lower()
        has_images = False
        if content_type == "pdf":
            has_images = has_images_in_pdf(source)
        elif content_type == "docx":
            has_images = has_images_in_docx(source)
        elif content_type in ["ppt", "pptx"]:
            has_images = has_images_in_pptx(source)
        elif content_type == "eml":
            has_images = has_images_in_eml(source)
        return DocumentMetadata(
            path=source,
            content_type=content_type,
            has_images=has_images,
            is_sharepoint=False
        )

@task.llm(
    model="gpt-4",
    task_id="determine_agents",
    result_type=List[AgentAssignment],
    system_prompt="""
    You are a document routing expert that analyzes document metadata and assigns appropriate processing agents.
    Consider the following agents:
    - File Parsing Agent: Handles PDF, PPT, PPTX, DOCX for text extraction
    - Email Agent: Processes email (EML) files
    - Transcript Agent: Processes text-heavy documents (VTT, DOCX)
    - Image Processing Agent: Handles documents with embedded images (has_images=True)

    Based on the metadata (content_type, has_images), assign the document to one or more agents with a priority (1-5, where 1 is highest). Rules:
    - PDF, PPT, PPTX: Assign to File Parsing Agent (priority 2)
    - EML: Assign to Email Agent (priority 2)
    - VTT, DOCX: Assign to Transcript Agent (priority 3)
    - If has_images=True: Assign to Image Processing Agent (priority 1)
    - DOCX can be assigned to both File Parsing and Transcript Agents
    Return a list of AgentAssignment objects.
    """
)
def determine_agents(metadata: DocumentMetadata) -> List[AgentAssignment]:
    pass

@tool
def extract_images(document_path: str) -> List[str]:
    """Extract images from the document and return their paths."""
    if document_path.endswith(".pdf"):
        from pdf2image import convert_from_path
        images = convert_from_path(document_path)
        image_paths = []
        for i, image in enumerate(images):
            image_path = f"/tmp/image_{i}.jpg"
            image.save(image_path, "JPEG")
            image_paths.append(image_path)
        return image_paths
    return []

@tool
def analyze_image(image_path: str) -> str:
    """Analyze the image using an AI model and return the result."""
    return "Image analysis result"

@task
def process_with_file_parsing_agent(assignment: AgentAssignment) -> dict:
    if assignment.agent_type == "file_parsing":
        logging.info(f"Processing {assignment.document.path} with File Parsing Agent (priority: {assignment.priority})")
        return {"status": "success", "agent": "file_parsing", "path": assignment.document.path}
    return None

@task
def process_with_email_agent(assignment: AgentAssignment) -> dict:
    if assignment.agent_type == "email":
        logging.info(f"Processing {assignment.document.path} with Email Agent (priority: {assignment.priority})")
        return {"status": "success", "agent": "email", "path": assignment.document.path}
    return None

@task
def process_with_transcript_agent(assignment: AgentAssignment) -> dict:
    if assignment.agent_type == "transcript":
        logging.info(f"Processing {assignment.document.path} with Transcript Agent (priority: {assignment.priority})")
        return {"status": "success", "agent": "transcript", "path": assignment.document.path}
    return None

@task.agent(model="gpt-4", tools=[extract_images, analyze_image])
def process_with_image_processing_agent(assignment: AgentAssignment) -> dict:
    if assignment.agent_type == "image_processing":
        logging.info(f"Processing {assignment.document.path} with Image Processing Agent (priority: {assignment.priority})")
        return {"status": "success", "agent": "image_processing", "path": assignment.document.path}
    return None

@dag(
    dag_id="intelligent_document_router",
    schedule=None,
    start_date=datetime(2025, 1, 1),
    catchup=False,
)
def document_router_dag():
    documents = ["/app/fdi/Documents/FRD.pptx"]
    sharepoint_urls = [
        "https://hexawareonline.sharepoint.com/teams/tensaiGPT-PROD-HR-Docs/_api/web/GetFileByServerRelativeUrl('/teams/tensaiGPT-PROD-HR-Docs/Shared%20Documents/document1.pdf')",
        "https://hexawareonline.sharepoint.com/teams/tensaiGPT-PROD-HR-Docs/Shared%20Documents/Forms/AllItems.aspx"
    ]

    local_metadata = validate_and_extract_metadata.expand(source=documents)
    sharepoint_metadata = validate_and_extract_metadata.expand(source=sharepoint_urls, is_sharepoint=True)
    all_metadata = local_metadata + sharepoint_metadata

    assignments = determine_agents.expand(metadata=all_metadata)

    @task
    def flatten_assignments(assignments_list):
        return [assignment for sublist in assignments_list for assignment in sublist]

    all_assignments = flatten_assignments(assignments)

    file_parsing_tasks = process_with_file_parsing_agent.expand(assignment=all_assignments)
    email_tasks = process_with_email_agent.expand(assignment=all_assignments)
    transcript_tasks = process_with_transcript_agent.expand(assignment=all_assignments)
    image_processing_tasks = process_with_image_processing_agent.expand(assignment=all_assignments)

    all_assignments >> [file_parsing_tasks, email_tasks, transcript_tasks, image_processing_tasks]

dag_instance = document_router_dag()
