import os
from typing import List, Dict, Any
from airflow import DAG
from airflow.decorators import dag, task
from datetime import datetime
from pydantic import BaseModel, Field
from pydantic_ai import Agent as PydanticAIAgent
from dotenv import load_dotenv
from airflow.utils.log.logging_mixin import LoggingMixin
from pdf2image import convert_from_path
from pptx import Presentation
from openpyxl import load_workbook
import json
import ast
# Load environment variables
load_dotenv()
api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise ValueError("OPENAI_API_KEY is required")
os.environ["OPENAI_API_KEY"] = api_key

# Logger
logger = LoggingMixin().log

# Tool function
def detect_images_in_document(path: str) -> bool:
    try:
        ext = os.path.splitext(path)[1].lower()
        if ext == ".pdf":
            images = convert_from_path(path)
            return len(images) > 0
        elif ext == ".pptx":
            prs = Presentation(path)
            return any(shape.shape_type == 13 for slide in prs.slides for shape in slide.shapes)
        elif ext == ".xlsx":
            wb = load_workbook(path)
            return any(sheet._images for sheet in wb)
        return False
    except Exception as e:
        logger.error(f"Error detecting images in {path}: {e}")
        return False

# System prompt
SYSTEM_PROMPT = """
You are a document classifier. Based on the input path (file path or URL), classify which agents should process the document.

- For '.pdf', '.pptx', '.xlsx': always classify as 'File Parsing Agent'
- For those same files('.pdf', '.pptx', '.xlsx'), if detect_images_in_document returns true, also classify as 'Image Processing Agent'
- For '.eml', '.msg': classify as 'Email Agent',
- For those same files('.eml', '.msg'), if detect_images_in_document returns true, also classify as 'Image Processing Agent'
- For '.vtt', '.txt': classify as 'Transcript Agent'
- For image extensions ('.jpg', '.png', '.gif'): classify as 'Image Processing Agent'
- For unknown or no extension: default to 'File Parsing Agent'
- If the path is a SharePoint URL (starts with 'http'), assume it may contain images

Always return a list of classification objects like:
[{"agent": "File Parsing Agent", "path": "<input>"}, {"agent": "Image Processing Agent", "path": "<input>"}]
"""

# Initialize Pydantic AI Agent (no output_model)
document_classifier_agent = PydanticAIAgent(
    model="gpt-4o",
    system_prompt=SYSTEM_PROMPT,
    tools=[detect_images_in_document],
)

@dag(
    dag_id="document_classifier_dag",
    schedule=None,
    start_date=datetime(2025, 1, 1),
    catchup=False,
    tags=["document_classification"],
)
def doc_classifier_dag():

    @task.agent(agent=document_classifier_agent)
    def classify(input_path: str) -> List[Dict[str, Any]]:
        logger.info(f"Running classifier for: {input_path}")
        result = document_classifier_agent.run_sync(input_path)
        output_data = result.data
        
        if isinstance(output_data, str):
            try:
                data = json.loads(output_data)
            except json.JSONDecodeError:
                try:
                    data = ast.literal_eval(output_data)
                except Exception:
                    logger.error(f"Failed to parse agent output: {output_data}")
                    raise
        else:
            data = output_data

        logger.info(f"Agent output: {data}")
        return data

    @task
    def flatten(results: List[List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        return [item for sublist in results for item in sublist]

    @task
    def process(assignment: Dict[str, Any]) -> Dict[str, Any]:
        logger.info(f"Assignment: Agent={assignment['agent']}, Path={assignment['path']}")
        return {**assignment, "status": "processed"}

    @task
    def summarize(results: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        logger.info("Final Summary:")
        for result in results:
            logger.info(result)
        return results

    # Inputs
    inputs = [
        "/app/fdi/Documents/FRD.pptx",
        "https://hexawareonline.sharepoint.com/:b:/t/tensaiGPT-PROD-HR-Docs/ET0W0clrClhBrA7ZLzCoOmEBHq0vg-rFuGuEwb40Weq8zQ?e=6BkU3C",
    ]

    classifications = classify.expand(input_path=inputs)
    flat = flatten(classifications)
    processed = process.expand(assignment=flat)
    summarize(processed)

doc_classifier_dag = doc_classifier_dag()
